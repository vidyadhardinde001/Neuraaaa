import os
import fitz
from pptx import Presentation
import docx

def get_file_content(file_path):
    _, file_extension = os.path.splitext(file_path)
    text = ""
    try:
        if file_extension.lower() == ".pdf":
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text()
        elif file_extension.lower() == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_extension.lower() == ".docx":
            document = docx.Document(file_path)
            for para in document.paragraphs:
                text += para.text + "\n"
        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
    except FileNotFoundError:
        return None, "Error: File not found. Please check the path and try again."
    except Exception as e:
        return None, f"Error processing file: {e}"
    if not text.strip():
        return None, "Error: Could not extract any text. The file might be empty or image-based."
    return text, None

import google.generativeai as genai

def process_file_with_gemini(content, api_key):
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('models/gemini-flash-latest')
        prompt = f"""
        Analyze the following file content.
        First, provide a summary of the file in a single, concise paragraph (3-4 sentences).
        Second, on a new line, suggest a new file name for this document. The file name should:
        - Be 3-5 words long.
        - Clearly describe the file's content.
        - Use hyphens instead of spaces.
        - Start with 'FileName:'.
        File Content:
        ---
        {content}
        ---
        """
        response = model.generate_content(prompt)
        full_response = response.text
        summary_lines = []
        suggested_name = "Not available."
        for line in full_response.splitlines():
            clean_line = line.strip()
            if clean_line.startswith("FileName:"):
                suggested_name = clean_line.replace("FileName:", "").strip()
            elif clean_line:
                summary_lines.append(clean_line)
        summary = "\n".join(summary_lines)
        return summary, suggested_name, None
    except Exception as e:
        return None, None, f"An error occurred with the Gemini API: {e}"
